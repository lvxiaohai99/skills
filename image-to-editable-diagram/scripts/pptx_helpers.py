#!/usr/bin/env python3
"""生成可编辑 PPTX 架构图的通用图元库。

封装 python-pptx 缺少直接 API 的几个痛点：
  - 中文东亚字体（a:ea）           -> _set_font
  - 虚线                            -> dash 参数（a:prstDash）
  - 连接符箭头                      -> arrow()（a:tailEnd）
  - 无阴影的干净矢量外观            -> shadow.inherit=False

坐标约定：虚拟画布单位 = 0.01 英寸。
构造时把幻灯片尺寸设为「画布单位/100」英寸，坐标即像素级直观。
先用 scripts/prepare_image.py 精读原图、量好各模块的相对位置，再套用本库摆放。

用法示例：
    from pptx_helpers import Diagram, BLUE, GREEN, L_BLUE
    d = Diagram(width_in=13.333, height_in=8.0)          # 16:9
    d.label(0, 8, 1333, 40, "标题", INK, size=20, bold=True, align="center")
    d.container(90, 55, 1200, 95, BLUE, dash=True)        # 预留域(虚线)
    b1 = d.box(185, 72, 180, 52, "云端平台", fill=L_BLUE, line=BLUE, bold=True)
    d.box(595, 72, 200, 52, "协议交互", sub="(设备认证/鉴权)", fill=L_BLUE, line=BLUE)
    d.arrow(365, 98, 595, 98, color=BLUE, label="下发", dash=False)
    d.save("out.pptx")
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ---------- 常用配色（可按图例增删） ----------
BLUE   = RGBColor(0x25, 0x63, 0xEB)
GREEN  = RGBColor(0x16, 0xA3, 0x4A)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
RED    = RGBColor(0xDC, 0x26, 0x26)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
INK    = RGBColor(0x1F, 0x29, 0x37)
GRAY   = RGBColor(0x6B, 0x72, 0x80)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
# 浅底（描边同色系）
L_BLUE  = RGBColor(0xEF, 0xF6, 0xFF)
L_GREEN = RGBColor(0xF0, 0xFD, 0xF4)
L_ORNG  = RGBColor(0xFF, 0xF7, 0xED)
L_RED   = RGBColor(0xFE, 0xF2, 0xF2)
L_PURP  = RGBColor(0xF5, 0xF3, 0xFF)
L_GRAY  = RGBColor(0xFA, 0xFA, 0xFA)

_ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}


class Diagram:
    """一页式架构图画布。所有坐标单位 = 0.01 英寸。"""

    def __init__(self, width_in=13.333, height_in=8.0, font="Microsoft YaHei",
                 mono="Consolas"):
        self.font, self.mono = font, mono
        self.prs = Presentation()
        self.prs.slide_width = Inches(width_in)
        self.prs.slide_height = Inches(height_in)
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # 空白

    # ---- 坐标换算 ----
    @staticmethod
    def U(v):
        return Inches(v / 100.0)

    # ---- 字体（含中文东亚字形） ----
    def _set_font(self, run, size, color, bold=False, name=None):
        name = name or self.font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = name
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn('a:ea'))
        if ea is None:
            ea = rPr.makeelement(qn('a:ea'), {})
            rPr.append(ea)
        ea.set('typeface', name)

    @staticmethod
    def _dash(shp):
        ln = shp.line._get_or_add_ln()
        d = ln.find(qn('a:prstDash'))
        if d is None:
            d = ln.makeelement(qn('a:prstDash'), {})
            ln.append(d)
        d.set('val', 'dash')

    # ---- 盒子（圆角矩形，可带副标题小字） ----
    def box(self, x, y, w, h, text, fill=WHITE, line=INK, fsize=11, bold=False,
            fcolor=INK, sub=None, line_w=1.5, dash=False, radius=True):
        shp = self.slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
            self.U(x), self.U(y), self.U(w), self.U(h))
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
        shp.shadow.inherit = False
        if dash:
            self._dash(shp)
        tf = shp.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Pt(3)
        tf.margin_top = tf.margin_bottom = Pt(2)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = text
        self._set_font(r, fsize, fcolor, bold)
        if sub:
            p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run(); r2.text = sub
            self._set_font(r2, fsize - 2, GRAY)
        return shp

    # ---- 分组容器（虚线域/实体域） ----
    def container(self, x, y, w, h, line, dash=True, fill=None, line_w=2.0):
        shp = self.slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, self.U(x), self.U(y), self.U(w), self.U(h))
        if fill is None:
            shp.fill.background()
        else:
            shp.fill.solid(); shp.fill.fore_color.rgb = fill
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
        shp.shadow.inherit = False
        if dash:
            self._dash(shp)
        return shp

    # ---- 独立文字标签 ----
    def label(self, x, y, w, h, text, color, size=12, bold=True, align="left"):
        tb = self.slide.shapes.add_textbox(self.U(x), self.U(y), self.U(w), self.U(h))
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = Pt(1)
        tf.margin_top = tf.margin_bottom = Pt(1)
        for i, line in enumerate(text.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = _ALIGN.get(align, PP_ALIGN.LEFT)
            r = p.add_run(); r.text = line
            self._set_font(r, size, color, bold)
        return tb

    # ---- 多行文本（目录树/说明列表），lines=[(文本, bold, color|None)] ----
    def multiline(self, x, y, w, h, lines, size=10, mono=False, color=INK):
        tb = self.slide.shapes.add_textbox(self.U(x), self.U(y), self.U(w), self.U(h))
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = Pt(2)
        for i, (txt, bold, col) in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = txt
            self._set_font(r, size, col or color, bold,
                           self.mono if mono else self.font)
        return tb

    # ---- 带箭头连接符（可虚线、可加线上标签） ----
    def arrow(self, x1, y1, x2, y2, color=GRAY, dash=False, width=1.75, label=None,
              label_size=8):
        cxn = self.slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, self.U(x1), self.U(y1), self.U(x2), self.U(y2))
        cxn.line.color.rgb = color; cxn.line.width = Pt(width)
        cxn.shadow.inherit = False
        ln = cxn.line._get_or_add_ln()
        tail = ln.makeelement(qn('a:tailEnd'), {})
        tail.set('type', 'triangle'); tail.set('w', 'med'); tail.set('len', 'med')
        ln.append(tail)
        if dash:
            d = ln.makeelement(qn('a:prstDash'), {}); d.set('val', 'dash'); ln.append(d)
        if label:
            mx, my = (x1 + x2) / 2, min(y1, y2) - 18
            self.label(mx - 45, my, 100, 16, label, color, size=label_size,
                       bold=True, align="center")
        return cxn

    # ---- 图例色样（一条短线 + 文字） ----
    def legend_item(self, x, y, text, color, dash=False, line_w=2.5,
                    swatch_len=36, text_size=9.5):
        ln = self.slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, self.U(x), self.U(y + 6),
            self.U(x + swatch_len), self.U(y + 6))
        ln.line.color.rgb = color; ln.line.width = Pt(line_w)
        ln.shadow.inherit = False
        if dash:
            e = ln.line._get_or_add_ln()
            d = e.makeelement(qn('a:prstDash'), {}); d.set('val', 'dash'); e.append(d)
        self.label(x + swatch_len + 6, y, 460, 16, text,
                   RGBColor(0x37, 0x41, 0x51), size=text_size, bold=False)

    def save(self, path):
        self.prs.save(path)
        return path
