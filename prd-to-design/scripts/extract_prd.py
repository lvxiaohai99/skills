#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""从 PRD 文件提取纯文本，供 prd-to-design 后续拆需求 / 写概要设计。

支持: pdf / docx / doc / pptx / ppt / xlsx / xls / md / txt / html / csv
用法: python extract_prd.py <文件路径> [更多文件...]
"""

from __future__ import annotations

import csv
import json
import os
import shutil
import subprocess
import sys
import tempfile
import traceback
from html.parser import HTMLParser
from typing import List, Tuple


def _force_utf8_stdio() -> None:
    """Windows 终端默认代码页常把中文打成乱码；强制 stdout/stderr 用 UTF-8。"""
    os.environ.setdefault("PYTHONIOENCODING", "utf-8")
    for stream in (sys.stdout, sys.stderr):
        try:
            stream.reconfigure(encoding="utf-8", errors="replace")
        except Exception:
            pass


_force_utf8_stdio()

NATIVE = {
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".xls",
    ".md",
    ".markdown",
    ".txt",
    ".html",
    ".htm",
    ".csv",
}
CONVERTIBLE = {".doc": "docx", ".ppt": "pptx"}
SUPPORTED = NATIVE | set(CONVERTIBLE)


def fail(message: str, code: int = 1) -> None:
    print(f"[extract_prd] ERROR: {message}", file=sys.stderr)
    sys.exit(code)


def read_text_file(path: str) -> str:
    raw = open(path, "rb").read()
    for enc in ("utf-8-sig", "utf-8", "gb18030", "gbk", "latin-1"):
        try:
            return raw.decode(enc)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace")


def _low_text(text: str) -> List[str]:
    return ["low_text"] if len(text.strip()) < 200 else []


def extract_pdf(path: str) -> Tuple[str, List[str]]:
    warnings: List[str] = []
    try:
        import pdfplumber
    except ImportError:
        fail("缺少 pdfplumber。请执行: pip install -r requirements.txt")
    chunks: List[str] = []
    with pdfplumber.open(path) as pdf:
        if not pdf.pages:
            warnings.append("empty_pdf")
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            table_lines: List[str] = []
            for table in page.extract_tables() or []:
                for row in table:
                    cells = [str(c).strip() if c is not None else "" for c in row]
                    if any(cells):
                        table_lines.append(" | ".join(cells))
            page_body = text.strip()
            if table_lines:
                page_body = (page_body + "\n" + "\n".join(table_lines)).strip()
            if page_body:
                chunks.append(f"--- 第 {i} 页 ---\n{page_body}")
    if not chunks:
        warnings.append("low_text")
        warnings.append("可能是扫描件 PDF，无法提取文字")
    return "\n\n".join(chunks), warnings


def extract_docx(path: str) -> Tuple[str, List[str]]:
    try:
        from docx import Document
    except ImportError:
        fail("缺少 python-docx。请执行: pip install -r requirements.txt")
    doc = Document(path)
    parts: List[str] = []
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    text = "\n".join(parts)
    return text, _low_text(text)


def _shape_texts(shape) -> List[str]:
    lines: List[str] = []
    if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
        t = (shape.text_frame.text or "").strip()
        if t:
            lines.append(t)
    if getattr(shape, "has_table", False) and shape.has_table:
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                lines.append(" | ".join(cells))
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                lines.extend(_shape_texts(child))
    except Exception:
        pass
    return lines


def extract_pptx(path: str) -> Tuple[str, List[str]]:
    try:
        from pptx import Presentation
    except ImportError:
        fail("缺少 python-pptx。请执行: pip install -r requirements.txt")
    prs = Presentation(path)
    chunks: List[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: List[str] = []
        for shape in slide.shapes:
            parts.extend(_shape_texts(shape))
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        body = "\n".join(parts)
        if notes:
            body = (body + "\n[备注]\n" + notes).strip()
        if body:
            chunks.append(f"--- 第 {i} 页 ---\n{body}")
    text = "\n\n".join(chunks)
    warnings = _low_text(text)
    if warnings:
        warnings.append("PPT 文本很少，可能以图为主；架构图需另补或人工说明")
    return text, warnings


def _sheet_to_text(rows: List[List[object]], sheet_name: str) -> str:
    lines = [f"--- 工作表: {sheet_name} ---"]
    for row in rows:
        cells = ["" if c is None else str(c).strip() for c in row]
        if any(cells):
            lines.append(" | ".join(cells))
    return "\n".join(lines)


def extract_xlsx(path: str) -> Tuple[str, List[str]]:
    try:
        import openpyxl
    except ImportError:
        fail("缺少 openpyxl。请执行: pip install -r requirements.txt")
    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    chunks = []
    for name in wb.sheetnames:
        ws = wb[name]
        rows = [list(row) for row in ws.iter_rows(values_only=True)]
        chunks.append(_sheet_to_text(rows, name))
    wb.close()
    text = "\n\n".join(chunks)
    return text, _low_text(text)


def extract_xls(path: str) -> Tuple[str, List[str]]:
    try:
        import xlrd
    except ImportError:
        fail("缺少 xlrd。请执行: pip install -r requirements.txt")
    book = xlrd.open_workbook(path)
    chunks = []
    for sheet in book.sheets():
        rows = [sheet.row_values(r) for r in range(sheet.nrows)]
        chunks.append(_sheet_to_text(rows, sheet.name))
    text = "\n\n".join(chunks)
    return text, _low_text(text)


class _HTMLText(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self._skip = False
        self.parts: List[str] = []

    def handle_starttag(self, tag: str, attrs) -> None:
        if tag in ("script", "style"):
            self._skip = True

    def handle_endtag(self, tag: str) -> None:
        if tag in ("script", "style"):
            self._skip = False
        if tag in ("p", "div", "br", "li", "tr", "h1", "h2", "h3", "h4"):
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if not self._skip and data.strip():
            self.parts.append(data.strip())


def extract_html(path: str) -> Tuple[str, List[str]]:
    parser = _HTMLText()
    parser.feed(read_text_file(path))
    text = " ".join(parser.parts)
    text = "\n".join(line.strip() for line in text.splitlines() if line.strip())
    return text, _low_text(text)


def extract_csv(path: str) -> Tuple[str, List[str]]:
    raw = read_text_file(path)
    rows = list(csv.reader(raw.splitlines()))
    lines = [" | ".join(c.strip() for c in row) for row in rows if any(x.strip() for x in row)]
    text = "\n".join(lines)
    return text, _low_text(text)


def _find_soffice() -> str | None:
    for name in ("soffice", "libreoffice"):
        found = shutil.which(name)
        if found:
            return found
    return None


def convert_with_soffice(path: str, to_ext: str) -> str:
    """把旧版 .doc / .ppt 转成 OOXML，返回转换后的临时文件路径（调用方负责删除目录）。"""
    soffice = _find_soffice()
    if not soffice:
        fail(
            f"当前是旧版 {os.path.splitext(path)[1]}，本机没有 LibreOffice。"
            f"请另存为 .{to_ext} 后重试。"
        )
    tmp = tempfile.mkdtemp(prefix="prd_extract_")
    cmd = [soffice, "--headless", "--convert-to", to_ext, "--outdir", tmp, os.path.abspath(path)]
    try:
        proc = subprocess.run(cmd, check=False, timeout=120, capture_output=True)
    except subprocess.TimeoutExpired:
        shutil.rmtree(tmp, ignore_errors=True)
        fail("LibreOffice 转换超时。请手动另存为 ." + to_ext + " 后重试。")
    if proc.returncode != 0:
        err = (proc.stderr or proc.stdout or b"").decode("utf-8", errors="replace")
        shutil.rmtree(tmp, ignore_errors=True)
        fail(f"LibreOffice 转换失败: {err.strip() or proc.returncode}")
    candidates = [
        os.path.join(tmp, f)
        for f in os.listdir(tmp)
        if f.lower().endswith("." + to_ext)
    ]
    if not candidates:
        shutil.rmtree(tmp, ignore_errors=True)
        fail("LibreOffice 未产出 ." + to_ext + "。请手动另存后重试。")
    return candidates[0]


def extract_native(path: str, ext: str) -> Tuple[str, List[str]]:
    if ext == ".pdf":
        return extract_pdf(path)
    if ext == ".docx":
        return extract_docx(path)
    if ext == ".pptx":
        return extract_pptx(path)
    if ext == ".xlsx":
        return extract_xlsx(path)
    if ext == ".xls":
        return extract_xls(path)
    if ext in {".html", ".htm"}:
        return extract_html(path)
    if ext == ".csv":
        return extract_csv(path)
    text = read_text_file(path)
    return text, _low_text(text)


def extract(path: str) -> Tuple[str, str, List[str]]:
    ext = os.path.splitext(path)[1].lower()
    if ext not in SUPPORTED:
        fail(f"不支持的格式: {ext}。支持: {', '.join(sorted(SUPPORTED))}")
    tmp_parent = None
    work_path, work_ext = path, ext
    warnings: List[str] = []
    try:
        if ext in CONVERTIBLE:
            target = CONVERTIBLE[ext]
            converted = convert_with_soffice(path, target)
            tmp_parent = os.path.dirname(converted)
            work_path, work_ext = converted, "." + target
            warnings.append(f"converted_{ext.lstrip('.')}_via_soffice")
        text, extra = extract_native(work_path, work_ext)
        warnings.extend(extra)
        return ext.lstrip("."), text.strip(), warnings
    finally:
        if tmp_parent:
            shutil.rmtree(tmp_parent, ignore_errors=True)


def extract_many(paths: List[str]) -> Tuple[str, dict]:
    file_metas = []
    bodies: List[str] = []
    all_warnings: List[str] = []
    for path in paths:
        if not os.path.isfile(path):
            fail(f"文件不存在: {path}")
        fmt, text, warnings = extract(path)
        abs_path = os.path.abspath(path)
        file_metas.append(
            {
                "file": abs_path,
                "format": fmt,
                "chars": len(text),
                "warnings": warnings,
            }
        )
        all_warnings.extend(warnings)
        header = f"===== 文件: {os.path.basename(path)} ====="
        bodies.append(header + "\n" + (text if text else "(无文本)"))
    combined = "\n\n".join(bodies)
    meta = {
        "files": file_metas,
        "chars": len(combined),
        "warnings": sorted(set(all_warnings)),
    }
    return combined, meta


def main() -> None:
    if len(sys.argv) < 2:
        fail("用法: python extract_prd.py <PRD文件> [更多文件...]")
    paths = sys.argv[1:]
    try:
        text, meta = extract_many(paths)
    except SystemExit:
        raise
    except Exception:
        traceback.print_exc()
        fail("解析失败，见上方堆栈。可改用 .docx / .md / .txt / .pptx / .xlsx 后重试。")
    print("===PRD_META===")
    print(json.dumps(meta, ensure_ascii=False, indent=2))
    print("===PRD_TEXT===")
    print(text if text else "(无文本)")


if __name__ == "__main__":
    main()
